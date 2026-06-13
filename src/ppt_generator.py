import re
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from src.utils import ensure_dir
from pathlib import Path

# 16:9 widescreen
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Visual style
ACCENT_ORANGE = RGBColor(0xF5, 0x82, 0x20)
TITLE_DARK = RGBColor(0x21, 0x21, 0x21)
BODY_GRAY = RGBColor(0x40, 0x40, 0x40)
FOOTER_GRAY = RGBColor(0xA6, 0xA6, 0xA6)

FONT_NAME = "Calibri"
FONT_NAME_CJK = "Microsoft YaHei"
FOOTER_TEXT = "AI Course Video Generator"
SECTION_LABEL = "KEY CONCEPTS"
SECTION_LABEL_CJK = "核心要点"
SUBTITLE_TEXT = "An AI-generated course explanation video"
SUBTITLE_TEXT_CJK = "由 AI 自动生成的课程讲解视频"

_CJK_RE = re.compile(r"[一-鿿]")


def _course_has_cjk(course_data: dict) -> bool:
    """True if any slide title or bullet contains Chinese characters."""
    if _CJK_RE.search(course_data.get("course_title", "")):
        return True
    return any(
        _CJK_RE.search(slide["title"]) or any(_CJK_RE.search(b) for b in slide["bullets"])
        for slide in course_data["slides"]
    )


def generate_ppt(course_data: dict, output_path: str) -> None:
    """
    Generate a 16:9 PPTX from the structured course data:
    one title slide followed by one content slide per course slide.
    Only slide titles and bullets are used (narration is for TTS later).
    Chinese slide text automatically switches to a CJK-capable font.
    """
    cjk = _course_has_cjk(course_data)
    font_name = FONT_NAME_CJK if cjk else FONT_NAME
    section_label = SECTION_LABEL_CJK if cjk else SECTION_LABEL
    subtitle_text = SUBTITLE_TEXT_CJK if cjk else SUBTITLE_TEXT

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    blank_layout = prs.slide_layouts[6]

    _add_title_slide(prs, blank_layout, course_data["course_title"], font_name, subtitle_text)

    total = len(course_data["slides"])
    for slide_data in course_data["slides"]:
        _add_content_slide(prs, blank_layout, slide_data, total, font_name, section_label, cjk)

    output_path = Path(output_path)
    ensure_dir(output_path.parent)
    prs.save(str(output_path))


def _add_title_slide(prs, layout, course_title: str, font_name: str, subtitle_text: str) -> None:
    slide = prs.slides.add_slide(layout)

    _add_accent_bar(slide, left=Inches(1.0), top=Inches(2.55), width=Inches(1.6))

    title_box = slide.shapes.add_textbox(
        Inches(1.0), Inches(2.8), Inches(11.3), Inches(1.6)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = course_title
    _style_run(run, size=40, bold=True, color=TITLE_DARK, font_name=font_name)

    subtitle_box = slide.shapes.add_textbox(
        Inches(1.0), Inches(4.45), Inches(11.3), Inches(0.6)
    )
    tf = subtitle_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = subtitle_text
    _style_run(run, size=18, bold=False, color=BODY_GRAY, font_name=font_name)

    footer_box = slide.shapes.add_textbox(
        Inches(1.0), Inches(6.9), Inches(11.3), Inches(0.4)
    )
    p = footer_box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = FOOTER_TEXT
    _style_run(run, size=10, bold=False, color=FOOTER_GRAY, font_name=font_name)


def _add_content_slide(prs, layout, slide_data: dict, total_slides: int,
                       font_name: str, section_label: str, cjk: bool) -> None:
    slide = prs.slides.add_slide(layout)

    _add_accent_bar(slide, left=Inches(0.8), top=Inches(0.7), width=Inches(1.2))

    label_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(0.85), Inches(4.0), Inches(0.35)
    )
    p = label_box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = section_label
    _style_run(run, size=11, bold=True, color=ACCENT_ORANGE, font_name=font_name)

    title_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(1.2), Inches(11.7), Inches(1.0)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = slide_data["title"]
    _style_run(run, size=30, bold=True, color=TITLE_DARK, font_name=font_name)

    bullets_box = slide.shapes.add_textbox(
        Inches(1.0), Inches(2.6), Inches(10.8), Inches(3.6)
    )
    tf = bullets_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(slide_data["bullets"]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(22)
        if cjk:
            p.line_spacing = 1.25  # extra leading for CJK readability

        marker = p.add_run()
        marker.text = "▪  "  # small square bullet marker
        _style_run(marker, size=20, bold=True, color=ACCENT_ORANGE, font_name=font_name)

        text = p.add_run()
        text.text = bullet
        _style_run(text, size=20, bold=False, color=BODY_GRAY, font_name=font_name)

    footer_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(6.95), Inches(6.0), Inches(0.4)
    )
    p = footer_box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = FOOTER_TEXT
    _style_run(run, size=10, bold=False, color=FOOTER_GRAY, font_name=font_name)

    number_box = slide.shapes.add_textbox(
        Inches(10.5), Inches(6.95), Inches(2.0), Inches(0.4)
    )
    p = number_box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"{slide_data['slide_index']:02d} / {total_slides:02d}"
    _style_run(run, size=10, bold=False, color=FOOTER_GRAY, font_name=font_name)


def _add_accent_bar(slide, left, top, width) -> None:
    """Add the orange accent line used at the top of each slide."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.06)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_ORANGE
    bar.line.fill.background()
    bar.shadow.inherit = False


def _style_run(run, size: int, bold: bool, color: RGBColor, font_name: str = FONT_NAME) -> None:
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
